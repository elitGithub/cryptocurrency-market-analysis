"""
PPTX Report Generator - Creates PowerPoint presentations.
"""
import json
import os
import logging
from typing import Dict, List, Optional
import pandas as pd
from .base_generator import ReportGenerator

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PPTXReportGenerator(ReportGenerator):
    """Generates PowerPoint presentations using python-pptx library."""

    def generate_node(
        self,
        symbol: str,
        signal_data: Dict,
        exchange_results: pd.DataFrame,
        suggestions: List[str],
        df: Optional[pd.DataFrame],
        chart_path: Optional[str],
        date_prefix: str
    ) -> bool:
        """Generate a PowerPoint presentation."""
        if not PPTX_AVAILABLE:
            logging.error("python-pptx not installed. Run: pip install python-pptx")
            return False

        logging.info("Generating PowerPoint presentation...")

        pptx_filename = f"{date_prefix}Crypto_Market_Analysis.pptx"
        pptx_path = os.path.join(self.output_dir, pptx_filename)

        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            signal = signal_data.get('signal', 'HOLD')
            confidence = signal_data.get('confidence', 'LOW')
            reasoning = signal_data.get('reasoning', 'N/A')
            price = signal_data.get('price', 0.0)
            rsi = signal_data.get('rsi', 50.0)

            # Slide 1: Title
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(8)
            height = Inches(2)

            title_box = slide.shapes.add_textbox(left, top, width, height)
            tf = title_box.text_frame
            tf.text = "Cryptocurrency Market Analysis"
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

            subtitle_box = slide.shapes.add_textbox(left, Inches(4.5), width, Inches(1))
            stf = subtitle_box.text_frame
            stf.text = symbol
            sp = stf.paragraphs[0]
            sp.alignment = PP_ALIGN.CENTER
            sp.font.size = Pt(32)
            sp.font.color.rgb = RGBColor(200, 200, 200)

            # Background gradient simulation
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(102, 126, 234)

            # Slide 2: Executive Summary
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            tf = title_box.text_frame
            tf.text = "Executive Summary"
            p = tf.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True

            # Signal box
            signal_colors = {'BUY': RGBColor(46, 204, 113), 'SELL': RGBColor(231, 76, 60), 'HOLD': RGBColor(243, 156, 18)}
            signal_color = signal_colors.get(signal, RGBColor(128, 128, 128))

            signal_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(1.5))
            signal_box.fill.solid()
            signal_box.fill.fore_color.rgb = signal_color
            stf = signal_box.text_frame
            stf.text = signal
            sp = stf.paragraphs[0]
            sp.alignment = PP_ALIGN.CENTER
            sp.font.size = Pt(48)
            sp.font.bold = True
            sp.font.color.rgb = RGBColor(255, 255, 255)

            # Details
            details_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2))
            dtf = details_box.text_frame
            dtf.text = f"Confidence: {confidence}\nPrice: ${price:,.2f} | RSI: {rsi:.1f}\n{reasoning}"
            for p in dtf.paragraphs:
                p.font.size = Pt(16)

            # Slide 3: Exchange Comparison
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            tf = title_box.text_frame
            tf.text = "Exchange Comparison"
            p = tf.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True

            if not exchange_results.empty:
                rows = len(exchange_results) + 1
                cols = 4
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                height = Inches(4)

                table = slide.shapes.add_table(rows, cols, left, top, width, height).table

                # Header
                table.cell(0, 0).text = 'Exchange'
                table.cell(0, 1).text = 'Total Pairs'
                table.cell(0, 2).text = 'USDT Pairs'
                table.cell(0, 3).text = 'Has OHLCV'

                for i in range(cols):
                    cell = table.cell(0, i)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(46, 64, 83)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(255, 255, 255)
                            run.font.bold = True

                # Data
                for idx, row in enumerate(exchange_results.itertuples(), 1):
                    table.cell(idx, 0).text = str(row.name)
                    table.cell(idx, 1).text = str(row.total_spot_pairs)
                    table.cell(idx, 2).text = str(row.usdt_quoted_pairs)
                    table.cell(idx, 3).text = 'Yes' if row.supports_fetchOHLCV else 'No'

            # Slide 4: Technical Analysis
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            tf = title_box.text_frame
            tf.text = "Key Takeaways"
            p = tf.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True

            text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
            tf = text_box.text_frame
            for suggestion in suggestions[:3]:
                p = tf.add_paragraph()
                p.text = suggestion
                p.level = 0
                p.font.size = Pt(18)

            # Slide 5: Chart
            if chart_path and os.path.exists(chart_path):
                slide = prs.slides.add_slide(prs.slide_layouts[5])

                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
                tf = title_box.text_frame
                tf.text = "Technical Analysis Chart"
                p = tf.paragraphs[0]
                p.font.size = Pt(32)
                p.font.bold = True

                slide.shapes.add_picture(chart_path, Inches(1), Inches(1.5), width=Inches(8))

            # Slide 6: Moving Averages
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            tf = title_box.text_frame
            tf.text = "Understanding Moving Averages"
            p = tf.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True

            text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf = text_box.text_frame

            p = tf.add_paragraph()
            p.text = "Golden Cross: Short MA > Long MA = Bullish"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(46, 204, 113)

            p = tf.add_paragraph()
            p.text = "Death Cross: Short MA < Long MA = Bearish"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(231, 76, 60)

            # Slide 7: RSI
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            tf = title_box.text_frame
            tf.text = "Understanding RSI"
            p = tf.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True

            text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf = text_box.text_frame

            p = tf.add_paragraph()
            p.text = "RSI > 70: Overbought (potential pullback)"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(231, 76, 60)

            p = tf.add_paragraph()
            p.text = "RSI 30-70: Neutral zone"
            p.font.size = Pt(20)

            p = tf.add_paragraph()
            p.text = "RSI < 30: Oversold (potential bounce)"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(46, 204, 113)

            p = tf.add_paragraph()
            p.text = f"\nCurrent RSI: {rsi:.1f}"
            p.font.size = Pt(28)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            # Slide 8: Disclaimer
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(102, 126, 234)

            title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf = title_box.text_frame

            p = tf.paragraphs[0]
            p.text = "Important Disclaimer"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

            p = tf.add_paragraph()
            p.text = "\nThis report is for educational purposes only.\nNOT financial advice. Cryptocurrency trading has substantial risk.\nAlways DYOR and never invest more than you can afford to lose."
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(255, 255, 255)

            prs.save(pptx_path)
            logging.info(f"✓ PowerPoint generated: {pptx_filename}")
            return True

        except Exception as e:
            logging.error(f"Error during PPTX generation: {e}")
            return False

    def generate(
        self,
        symbol: str,
        signal_data: Dict,
        exchange_results: pd.DataFrame,
        suggestions: List[str],
        df: Optional[pd.DataFrame],
        chart_path: Optional[str],
        date_prefix: str
    ) -> bool:
        """Generate a PowerPoint presentation."""
        logging.info("Generating PowerPoint presentation...")

        # Prepare output path
        pptx_filename = f"{date_prefix}Crypto_Market_Analysis.pptx"
        pptx_path = os.path.abspath(os.path.join(self.output_dir, pptx_filename))

        # Create HTML slides
        html_files = self._create_html_slides(
            symbol, signal_data, exchange_results, suggestions, df
        )

        # Create JavaScript code
        js_code = self._generate_js_code(
            exchange_results, chart_path, pptx_path, len(html_files)
        )

        # Write script
        script_path = os.path.join(os.getcwd(), 'generate_pptx_temp.js')
        temp_files = [script_path] + html_files

        try:
            with open(script_path, 'w', encoding='utf-8') as f:
                f.write(js_code)

            # Copy chart if it exists
            if chart_path and os.path.exists(chart_path):
                import shutil
                chart_copy = os.path.join(os.getcwd(), os.path.basename(chart_path))
                if os.path.abspath(chart_path) != chart_copy:
                    shutil.copy(chart_path, chart_copy)
                    temp_files.append(chart_copy)

            # Run the script
            success = self._run_node_script(script_path, timeout=60)

            # Cleanup
            self._cleanup_temp_files(temp_files)

            if success and os.path.exists(pptx_path):
                logging.info(f"✓ PowerPoint generated: {pptx_filename}")
                return True
            else:
                logging.error("PowerPoint generation failed")
                return False

        except Exception as e:
            logging.error(f"Error during PPTX generation: {e}")
            self._cleanup_temp_files(temp_files)
            return False

    def _create_html_slides(
        self,
        symbol: str,
        signal_data: Dict,
        exchange_results: pd.DataFrame,
        suggestions: List[str],
        df: Optional[pd.DataFrame]
    ) -> List[str]:
        """Create HTML files for each slide."""

        signal = signal_data.get('signal', 'HOLD')
        confidence = signal_data.get('confidence', 'LOW')
        reasoning = signal_data.get('reasoning', 'N/A')
        price = signal_data.get('price', 0.0)
        rsi = signal_data.get('rsi', 50.0)

        signal_colors = {'BUY': '#2ECC71', 'SELL': '#E74C3C', 'HOLD': '#F39C12'}
        signal_color = signal_colors.get(signal, '#808080')

        # Trend determination
        trend_text = "N/A"
        if df is not None and not df.empty and 'SMA_short' in df.columns and 'SMA_long' in df.columns:
            last = df.iloc[-1]
            if pd.notna(last['SMA_short']) and pd.notna(last['SMA_long']):
                trend_text = "Bullish" if last['SMA_short'] > last['SMA_long'] else "Bearish"

        css = """
:root { --primary: #2E4053; --success: #2ECC71; --danger: #E74C3C; --warning: #F39C12; }
body { font-family: Arial, sans-serif; margin: 0; padding: 40px; width: 920px; height: 500px; }
h1 { font-size: 48px; font-weight: 700; margin: 0 0 20px 0; color: var(--primary); }
h2 { font-size: 36px; font-weight: 600; margin: 0 0 15px 0; color: var(--primary); }
h3 { font-size: 24px; font-weight: 600; margin: 0 0 10px 0; }
p { font-size: 18px; line-height: 1.6; margin: 0 0 15px 0; }
.center { text-align: center; }
"""

        html_files = []

        # Slide 1: Title
        slide1 = f"""<!DOCTYPE html>
<html><head><style>{css}</style></head>
<body class="center" style="background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding-top: 150px;">
    <h1 style="color: white; font-size: 60px;">Cryptocurrency Market Analysis</h1>
    <h2 style="color: white; font-size: 48px;">{symbol}</h2>
    <p style="font-size: 20px; margin-top: 30px;">Comprehensive Technical Analysis Report</p>
</body></html>"""

        # Slide 2: Executive Summary
        slide2 = f"""<!DOCTYPE html>
<html><head><style>{css}</style></head>
<body style="background: white;">
    <h2>Executive Summary</h2>
    <div style="background: {signal_color}; color: white; padding: 30px; border-radius: 12px; margin: 20px 0;">
        <h1 style="color: white; font-size: 56px;">{signal}</h1>
        <p style="font-size: 20px;">Confidence: {confidence}</p>
    </div>
    <p><strong>Price:</strong> ${price:,.2f} | <strong>RSI:</strong> {rsi:.1f} | <strong>Trend:</strong> {trend_text}</p>
    <p style="margin-top: 20px;">{reasoning}</p>
</body></html>"""

        # Slide 3: Exchange Comparison with placeholder
        slide3 = f"""<!DOCTYPE html>
<html><head><style>{css}</style></head>
<body style="background: white;">
    <h2>Exchange Comparison</h2>
    <p>Analysis of multiple cryptocurrency exchanges:</p>
    <div class="placeholder" style="width: 100%; height: 300px; background: #f0f0f0;"></div>
</body></html>"""

        # Slide 4: Technical Analysis Chart with placeholder
        slide4 = f"""<!DOCTYPE html>
<html><head><style>{css}</style></head>
<body style="background: white;">
    <h2>Technical Analysis Chart</h2>
    <div class="placeholder" style="width: 100%; height: 400px; background: #f0f0f0;"></div>
</body></html>"""

        # Slide 5: Key Takeaways
        bullets = '<br>'.join([f"• {s[:100]}..." if len(s) > 100 else f"• {s}" for s in suggestions[:3]])
        slide5 = f"""<!DOCTYPE html>
<html><head><style>{css}</style></head>
<body style="background: white;">
    <h2>Key Takeaways</h2>
    <div style="margin-top: 30px; font-size: 20px; line-height: 2;">
        {bullets}
    </div>
</body></html>"""

        # Slide 6: Understanding Moving Averages
        slide6 = f"""<!DOCTYPE html>
<html><head><style>{css}</style></head>
<body style="background: white;">
    <h2>Understanding Moving Averages</h2>
    <p><strong style="color: var(--success);">Golden Cross:</strong> Short MA crosses above Long MA = Bullish</p>
    <p><strong style="color: var(--danger);">Death Cross:</strong> Short MA crosses below Long MA = Bearish</p>
    <p style="margin-top: 30px;">Moving averages smooth price data to identify trends over time.</p>
</body></html>"""

        # Slide 7: Understanding RSI
        slide7 = f"""<!DOCTYPE html>
<html><head><style>{css}</style></head>
<body style="background: white;">
    <h2>Understanding RSI</h2>
    <p style="color: var(--danger);"><strong>RSI &gt; 70:</strong> Overbought (potential pullback)</p>
    <p style="color: var(--warning);"><strong>RSI 30-70:</strong> Neutral zone</p>
    <p style="color: var(--success);"><strong>RSI &lt; 30:</strong> Oversold (potential bounce)</p>
    <p style="margin-top: 30px; font-size: 36px; text-align: center;"><strong>Current RSI: {rsi:.1f}</strong></p>
</body></html>"""

        # Slide 8: Disclaimer
        slide8 = """<!DOCTYPE html>
<html><head><style>{css}</style></head>
<body style="background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white;">
    <h2 style="color: white;">Important Disclaimer</h2>
    <p style="font-size: 20px; line-height: 1.8;">
        This report is for <strong>educational purposes only</strong>. 
        NOT financial advice. Cryptocurrency trading carries substantial risk.
        Always DYOR and never invest more than you can afford to lose.
    </p>
</body></html>""".replace('{css}', css)

        # Write HTML files
        slides = [slide1, slide2, slide3, slide4, slide5, slide6, slide7, slide8]
        for i, content in enumerate(slides, 1):
            filename = f'slide{i}.html'
            filepath = os.path.join(os.getcwd(), filename)
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write(content)
            html_files.append(filepath)

        return html_files

    def _generate_js_code(
        self,
        exchange_results: pd.DataFrame,
        chart_path: Optional[str],
        output_path: str,
        num_slides: int
    ) -> str:
        """Generate Node.js code for creating PowerPoint."""

        exchanges_json = json.dumps(
            exchange_results.to_dict('records') if isinstance(exchange_results, pd.DataFrame) else []
        )

        chart_filename = os.path.basename(chart_path) if chart_path and os.path.exists(chart_path) else None

        js_code = f"""
const pptxgen = require("pptxgenjs");
const {{ html2pptx }} = require("@ant/html2pptx");

(async () => {{
    try {{
        const pptx = new pptxgen();
        pptx.layout = "LAYOUT_16x9";
        pptx.author = "Crypto Trading Bot";
        pptx.title = "Market Analysis Report";
        
        const exchangeData = {exchanges_json};
        
        // Add slides 1-2
        await html2pptx("slide1.html", pptx);
        await html2pptx("slide2.html", pptx);
        
        // Add slide 3 with exchange table
        const {{ slide: slide3, placeholders: p3 }} = await html2pptx("slide3.html", pptx);
        
        const tableData = [[
            {{'text':'Exchange','options':{{'fill':{{'color':'2E4053'}},'color':'FFFFFF','bold':true}}}},
            {{'text':'Total Pairs','options':{{'fill':{{'color':'2E4053'}},'color':'FFFFFF','bold':true}}}},
            {{'text':'USDT Pairs','options':{{'fill':{{'color':'2E4053'}},'color':'FFFFFF','bold':true}}}},
            {{'text':'Has OHLCV','options':{{'fill':{{'color':'2E4053'}},'color':'FFFFFF','bold':true}}}}
        ]];
        
        exchangeData.forEach(ex => {{
            tableData.push([
                ex.name || 'N/A',
                String(ex.total_spot_pairs ?? 'N/A'),
                String(ex.usdt_quoted_pairs ?? 'N/A'),
                ex.supports_fetchOHLCV ? 'Yes' : 'No'
            ]);
        }});
        
        if (p3 && p3.length > 0) {{
            slide3.addTable(tableData, {{
                ...p3[0],
                border: {{pt: 1, color: 'CCCCCC'}},
                fontSize: 16,
                align: 'center',
                valign: 'middle'
            }});
        }}
        
        // Add slide 4 with chart
        const {{ slide: slide4, placeholders: p4 }} = await html2pptx("slide4.html", pptx);
        
        {"if (p4 && p4.length > 0) {" if chart_filename else "// No chart available"}
            {"slide4.addImage({" if chart_filename else ""}
                {"path: '" + (chart_filename or '') + "'," if chart_filename else ""}
                {"...p4[0]" if chart_filename else ""}
            {"});" if chart_filename else ""}
        {"}" if chart_filename else ""}
        
        // Add remaining slides
        await html2pptx("slide5.html", pptx);
        await html2pptx("slide6.html", pptx);
        await html2pptx("slide7.html", pptx);
        await html2pptx("slide8.html", pptx);
        
        await pptx.writeFile({json.dumps(output_path)});
        console.log("✓ PowerPoint generated successfully");
        
    }} catch (err) {{
        console.error("✗ Error generating PowerPoint:", err.message, err.stack);
        process.exit(1);
    }}
}})();
"""

        return js_code